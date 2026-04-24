"""Extract plain text from uploaded files."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader


def extract_text(filename: str, raw: bytes) -> str:
    """Return UTF-8 text for supported formats; raises ValueError if unsupported."""
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        return _pdf(raw)
    if ext in {".docx"}:
        return _docx(raw)
    if ext in {".pptx"}:
        return _pptx(raw)
    if ext in {".xlsx", ".xlsm"}:
        return _xlsx(raw)
    if ext in {".html", ".htm"}:
        return _html(raw)
    if ext in {".txt", ".md", ".csv"}:
        return raw.decode("utf-8", errors="replace")
    raise ValueError(f"Unsupported extension: {ext}")


def chunk_text(text: str, size: int, overlap: int) -> list[str]:
    """Character-based overlapping chunks."""
    text = text.strip()
    if not text:
        return []
    if size <= 0:
        return [text]
    chunks: list[str] = []
    start = 0
    n = len(text)
    while start < n:
        end = min(start + size, n)
        piece = text[start:end].strip()
        if piece:
            chunks.append(piece)
        if end >= n:
            break
        start = max(end - overlap, start + 1)
    return chunks


def _pdf(raw: bytes) -> str:
    reader = PdfReader(BytesIO(raw))
    parts: list[str] = []
    for page in reader.pages:
        t = page.extract_text() or ""
        parts.append(t)
    return "\n\n".join(parts).strip()


def _docx(raw: bytes) -> str:
    doc = DocxDocument(BytesIO(raw))
    return "\n".join(p.text for p in doc.paragraphs if p.text).strip()


def _pptx(raw: bytes) -> str:
    prs = Presentation(BytesIO(raw))
    lines: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                t = paragraph.text.strip()
                if t:
                    lines.append(t)
    return "\n".join(lines).strip()


def _xlsx(raw: bytes) -> str:
    wb = load_workbook(filename=BytesIO(raw), read_only=True, data_only=True)
    lines: list[str] = []
    for sheet in wb.worksheets:
        lines.append(f"## {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                lines.append("\t".join(cells))
    return "\n".join(lines).strip()


def _html(raw: bytes) -> str:
    soup = BeautifulSoup(raw, "lxml")
    return soup.get_text("\n", strip=True)
